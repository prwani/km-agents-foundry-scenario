from __future__ import annotations

import json
import os
from pathlib import Path
import re
from typing import Annotated

from agent_framework import tool
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from km_agents.contracts import CaseStudyContent, CaseStudyRequest
from km_agents.grounding import EVIDENCE_PENDING, text_is_evidence_backed
from km_agents.pptx_skill.generation import generate_case_study


_SUPPORTED_SOURCE_EXTENSIONS = frozenset({".docx", ".pdf", ".pptx", ".xlsx"})


def _workspace() -> Path:
    return Path(os.getenv("AGENT_WORKSPACE_ROOT", str(Path.home() / "data"))).resolve()


def _workspace_file(relative_path: str, directory: str, suffixes: frozenset[str]) -> Path:
    if not relative_path or Path(relative_path).is_absolute():
        raise ValueError("Workspace file path must be relative")

    requested_path = Path(relative_path)
    if requested_path.parts and requested_path.parts[0].casefold() == directory.casefold():
        requested_path = Path(*requested_path.parts[1:])

    root = (_workspace() / directory).resolve()
    candidate = (root / requested_path).resolve()
    if root not in candidate.parents or candidate.suffix.lower() not in suffixes:
        raise ValueError(f"Path must name a supported file under {directory}/")
    if not candidate.is_file():
        raise FileNotFoundError(f"Workspace file is missing: {relative_path}")
    return candidate


def _docx_text(path: Path) -> str:
    document = Document(path)
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _pdf_text(path: Path) -> str:
    return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages).strip()


def _pptx_text(path: Path) -> str:
    presentation = Presentation(path)
    return "\n".join(
        shape.text.strip()
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    )


def _xlsx_text(path: Path) -> str:
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        rows: list[str] = []
        for worksheet in workbook.worksheets:
            rows.append(f"[{worksheet.title}]")
            for row in worksheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value not in (None, "")]
                if values:
                    rows.append(" | ".join(values))
        return "\n".join(rows)
    finally:
        workbook.close()


def extract_source_text(relative_path: str) -> str:
    path = _workspace_file(relative_path, "input", _SUPPORTED_SOURCE_EXTENSIONS)
    extractors = {
        ".docx": _docx_text,
        ".pdf": _pdf_text,
        ".pptx": _pptx_text,
        ".xlsx": _xlsx_text,
    }
    text = extractors[path.suffix.lower()](path).strip()
    if not text:
        raise ValueError(f"No readable text was extracted from {relative_path}")
    return text


def _evidence_backed_content(content: CaseStudyContent) -> tuple[CaseStudyContent, int]:
    evidence = "\n".join(extract_source_text(path) for path in content.provenance_files)
    suppressed_count = 0

    def grounded(value: str) -> str:
        nonlocal suppressed_count
        if text_is_evidence_backed(value, evidence):
            return value
        suppressed_count += 1
        return EVIDENCE_PENDING

    return (
        content.model_copy(
            update={
                "title": grounded(content.title),
                "challenge": grounded(content.challenge),
                "solution_overview": grounded(content.solution_overview),
                "architecture_components": [grounded(value) for value in content.architecture_components],
                "implementation_steps": [grounded(value) for value in content.implementation_steps],
                "measurable_outcomes": [grounded(value) for value in content.measurable_outcomes],
                "customer_quote": grounded(content.customer_quote),
                "next_steps": [grounded(value) for value in content.next_steps],
            }
        ),
        suppressed_count,
    )


def _redact_unapproved_customer_name(
    content: CaseStudyContent,
    request: CaseStudyRequest,
) -> tuple[CaseStudyContent, int]:
    if request.customer_name_approved_for_external_use:
        return content, 0

    pattern = re.compile(re.escape(request.customer_name), re.IGNORECASE)
    replacements = 0

    def redact(value: str) -> str:
        nonlocal replacements
        replacement, count = pattern.subn(request.display_customer_name, value)
        replacements += count
        return replacement

    redacted = content.model_dump()
    for key, value in redacted.items():
        if isinstance(value, str):
            redacted[key] = redact(value)
        elif isinstance(value, list):
            redacted[key] = [redact(item) for item in value]
    return CaseStudyContent.model_validate(redacted), replacements


@tool(
    name="extract_uploaded_evidence",
    description=(
        "Extract text from an uploaded DOCX, PDF, PPTX, or XLSX under input/. "
        "Accepts either a path relative to input/ or an input/-prefixed workspace-relative path."
    ),
    approval_mode="never_require",
)
def extract_uploaded_evidence_tool(relative_path: str) -> str:
    return json.dumps({"path": relative_path, "text": extract_source_text(relative_path)})


@tool(
    name="generate_case_study_deck",
    description=(
        "Generate a PPTX from validated CaseStudyContent using the canonical template. "
        "Use the content object schema exactly; it contains all required source-grounded "
        "case-study fields. Every content field must be a direct source statement; unsupported "
        "content is replaced with Evidence pending. Pass the original CaseStudyRequest JSON so "
        "unapproved customer names can be redacted. Write only a relative .pptx filename under output/."
    ),
    approval_mode="never_require",
)
def create_case_study_deck_tool(
    content: CaseStudyContent,
    output_filename: Annotated[str, "Relative PPTX filename to write under output/."],
    request_json: Annotated[str, "Original CaseStudyRequest JSON."],
) -> str:
    if not output_filename or Path(output_filename).is_absolute():
        raise ValueError("Output filename must be a relative .pptx path")
    workspace = _workspace()
    output_root = (workspace / "output").resolve()
    output_path = (output_root / output_filename).resolve()
    if output_root not in output_path.parents or output_path.suffix.lower() != ".pptx":
        raise ValueError("Output filename must be a .pptx path under output/")

    request = CaseStudyRequest.model_validate_json(request_json)
    normalized_content, suppressed_outcomes = _evidence_backed_content(CaseStudyContent.model_validate(content))
    normalized_content, redacted_customer_name_occurrences = _redact_unapproved_customer_name(
        normalized_content,
        request,
    )
    template_path = Path(
        os.getenv("CASE_STUDY_TEMPLATE_PATH", "assets/templates/contoso-case-study-template.pptx")
    ).resolve()
    report = generate_case_study(template_path, output_path, normalized_content)
    return json.dumps(
        {
            "deck_path": f"output/{output_filename}",
            "report": report,
            "suppressed_unsupported_content_fields": suppressed_outcomes,
            "redacted_customer_name_occurrences": redacted_customer_name_occurrences,
        }
    )
