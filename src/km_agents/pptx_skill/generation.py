from __future__ import annotations

import hashlib
import json
from pathlib import Path

from pptx import Presentation

from km_agents.contracts import CaseStudyContent
from km_agents.grounding import EVIDENCE_PENDING


def _replace_text(shape: object, value: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise ValueError(f"Editable shape has no text frame: {shape.name}")
    frame = shape.text_frame
    paragraph = frame.paragraphs[0]
    if not paragraph.runs:
        paragraph.text = value
        return

    # Keep the placeholder's run formatting, including its foreground color.
    paragraph.runs[0].text = value
    for run in paragraph.runs[1:]:
        run.text = ""
    for extra_paragraph in frame.paragraphs[1:]:
        for run in extra_paragraph.runs:
            run.text = ""


def _template_items(items: list[str], count: int) -> list[str]:
    return (items[:count] + [EVIDENCE_PENDING] * count)[:count]


def generate_case_study(
    template_path: Path,
    output_path: Path,
    content: CaseStudyContent,
    report_path: Path | None = None,
) -> dict[str, object]:
    if not template_path.is_file():
        raise FileNotFoundError(f"Template does not exist: {template_path}")
    presentation = Presentation(template_path)
    replacements = {
        "editable:s1:customer-name": content.customer_display_name,
        "editable:s1:title": content.title,
        "editable:s1:subtitle": "A synthetic customer success story",
        "editable:s2:challenge": content.challenge,
        "editable:s2:context": "Evidence-backed modernization opportunity",
        "editable:s3:solution": content.solution_overview,
        "editable:s3:solution-pillars": "   •   ".join(
            _template_items(content.architecture_components, 3)
        ),
        "editable:s7:quote": content.customer_quote,
        "editable:s7:quote-attribution": f"{content.customer_display_name} stakeholder",
        "editable:s8:next-steps": "\n\n".join(
            f"{index}. {step}"
            for index, step in enumerate(_template_items(content.next_steps, 2), start=1)
        ),
    }
    architecture_groups = [[], [], []]
    for index, component in enumerate(_template_items(content.architecture_components, 3)):
        architecture_groups[index % 3].append(component)
    for index, components in enumerate(architecture_groups, start=1):
        replacements[f"editable:s4:architecture-column-{index}"] = "\n\n".join(components)
    steps = _template_items(content.implementation_steps, 4)
    for index in range(1, 5):
        replacements[f"editable:s5:implementation-step-{index}"] = steps[index - 1]
    outcomes = _template_items(content.measurable_outcomes, 3)
    for index in range(1, 4):
        replacements[f"editable:s6:outcome-{index}"] = outcomes[index - 1]
    replaced: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.name in replacements:
                _replace_text(shape, replacements[shape.name])
                replaced.append(shape.name)

    missing = sorted(set(replacements) - set(replaced))
    if missing:
        raise ValueError(f"Template is missing required editable shapes: {', '.join(missing)}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    report: dict[str, object] = {
        "output_sha256": hashlib.sha256(output_path.read_bytes()).hexdigest(),
        "output_size_bytes": output_path.stat().st_size,
        "replaced_shapes": sorted(replaced),
        "provenance_files": content.provenance_files,
    }
    if report_path:
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(report, indent=2), encoding="utf-8")
    return report
