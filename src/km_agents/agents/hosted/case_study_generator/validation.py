from __future__ import annotations

import json
import os
from pathlib import Path
import re

from agent_framework import tool
from pptx import Presentation

from km_agents.agents.hosted.case_study_generator.operations import extract_source_text
from km_agents.contracts import (
    CaseStudyRequest,
    FindingCode,
    FindingSeverity,
    ValidationFinding,
    ValidationResult,
)
from km_agents.grounding import EVIDENCE_PENDING, text_is_evidence_backed
from km_agents.pptx_skill.validation import validate_presentation

_EMU_PER_INCH = 914400
_SAFE_MARGIN_INCHES = 0.55


def workspace_deck_path(relative_path: str) -> Path:
    if not relative_path or Path(relative_path).is_absolute():
        raise ValueError("Deck path must be relative to the hosted session workspace")
    workspace = Path(os.getenv("AGENT_WORKSPACE_ROOT", str(Path.home() / "data"))).resolve()
    candidate = (workspace / relative_path).resolve()
    if workspace not in candidate.parents or candidate.suffix.lower() != ".pptx":
        raise ValueError("Deck path must resolve to a PPTX inside the hosted session workspace")
    if not candidate.is_file():
        raise FileNotFoundError(f"Deck is missing from the hosted session workspace: {relative_path}")
    return candidate


def _editable_content_segments(shape: object) -> list[str]:
    text = shape.text.strip()
    if not text:
        return []
    return [
        re.sub(r"^\d+\.\s*", "", segment).strip()
        for segment in re.split(r"[\n•]+|\s(?=\d+\.\s)", text)
        if segment.strip()
    ]


def _is_template_boilerplate(shape_name: str, value: str, request: CaseStudyRequest) -> bool:
    if shape_name == "editable:s1:customer-name":
        return value == request.display_customer_name
    if shape_name == "editable:s1:subtitle":
        return value == "A synthetic customer success story"
    if shape_name == "editable:s2:context":
        return value == "Evidence-backed modernization opportunity"
    return shape_name == "editable:s7:quote-attribution" and value == (
        f"{request.display_customer_name} stakeholder"
    )


def _source_form(value: str, request: CaseStudyRequest) -> str:
    if request.customer_name_approved_for_external_use:
        return value
    return re.sub(
        re.escape(request.display_customer_name),
        request.customer_name,
        value,
        flags=re.IGNORECASE,
    )


def _template_path() -> Path:
    return Path(
        os.getenv("CASE_STUDY_TEMPLATE_PATH", "assets/templates/contoso-case-study-template.pptx")
    )


def _font_style(shape: object) -> tuple[object, ...] | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    paragraph = shape.text_frame.paragraphs[0]
    if not paragraph.runs:
        return None
    font = paragraph.runs[0].font
    color = font.color
    try:
        rgb = color.rgb
    except AttributeError:
        rgb = None
    return (
        font.name,
        font.size,
        font.bold,
        font.italic,
        color.type,
        rgb,
    )


def _visual_finding(
    message: str,
    *,
    slide_number: int,
    shape_name: str | None = None,
    evidence: dict[str, object] | None = None,
) -> ValidationFinding:
    return ValidationFinding(
        code=FindingCode.VISUAL_BRAND_VIOLATION,
        severity=FindingSeverity.ERROR,
        message=message,
        slide_number=slide_number,
        shape_name=shape_name,
        evidence=evidence or {},
    )


def _estimated_line_count(shape: object) -> int:
    frame = shape.text_frame
    font_size = _font_style(shape)[1] if _font_style(shape) else None
    if font_size is None:
        return 0
    characters_per_line = max(
        1,
        int((shape.width / _EMU_PER_INCH * 72) / (font_size.pt * 0.52)),
    )
    lines = 0
    for paragraph in frame.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        lines += max(1, (len(text) + characters_per_line - 1) // characters_per_line)
    return lines


def _visual_brand_findings(deck_path: Path) -> list[ValidationFinding]:
    template_path = _template_path()
    if not template_path.is_file():
        return [
            ValidationFinding(
                code=FindingCode.INCONCLUSIVE,
                severity=FindingSeverity.ERROR,
                message=f"Canonical template is unavailable for visual QA: {template_path}",
            )
        ]

    template = Presentation(template_path)
    candidate = Presentation(deck_path)
    findings: list[ValidationFinding] = []
    for slide_number, (template_slide, candidate_slide) in enumerate(
        zip(template.slides, candidate.slides, strict=False),
        start=1,
    ):
        expected_shapes = {shape.name: shape for shape in template_slide.shapes}
        actual_shapes = {shape.name: shape for shape in candidate_slide.shapes}
        unexpected_shapes = sorted(set(actual_shapes) - set(expected_shapes))
        for shape_name in unexpected_shapes:
            findings.append(
                _visual_finding(
                    "Deck contains an unapproved visual element",
                    slide_number=slide_number,
                    shape_name=shape_name,
                )
            )

        for shape_name, expected in expected_shapes.items():
            if not shape_name.startswith("editable:"):
                continue
            actual = actual_shapes.get(shape_name)
            if actual is None or not getattr(actual, "has_text_frame", False):
                continue

            expected_geometry = (expected.left, expected.top, expected.width, expected.height)
            actual_geometry = (actual.left, actual.top, actual.width, actual.height)
            if actual_geometry != expected_geometry:
                findings.append(
                    _visual_finding(
                        "Editable content region no longer matches the approved layout",
                        slide_number=slide_number,
                        shape_name=shape_name,
                    )
                )
                continue

            if _font_style(actual) != _font_style(expected):
                findings.append(
                    _visual_finding(
                        "Editable content no longer preserves the approved typography or contrast",
                        slide_number=slide_number,
                        shape_name=shape_name,
                    )
                )

            left = actual.left / _EMU_PER_INCH
            top = actual.top / _EMU_PER_INCH
            right = (actual.left + actual.width) / _EMU_PER_INCH
            bottom = (actual.top + actual.height) / _EMU_PER_INCH
            slide_width = candidate.slide_width / _EMU_PER_INCH
            slide_height = candidate.slide_height / _EMU_PER_INCH
            if (
                left < _SAFE_MARGIN_INCHES
                or top < _SAFE_MARGIN_INCHES
                or right > slide_width - _SAFE_MARGIN_INCHES
                or bottom > slide_height - _SAFE_MARGIN_INCHES
            ):
                findings.append(
                    _visual_finding(
                        "Editable content region exceeds the 0.55-inch safe margin",
                        slide_number=slide_number,
                        shape_name=shape_name,
                    )
                )

            font_size = _font_style(actual)[1] if _font_style(actual) else None
            if font_size:
                available_lines = int(
                    (actual.height / _EMU_PER_INCH * 72) / (font_size.pt * 1.3)
                )
                if _estimated_line_count(actual) > available_lines:
                    findings.append(
                        _visual_finding(
                            "Editable text is likely to overflow its approved visual region",
                            slide_number=slide_number,
                            shape_name=shape_name,
                        )
                    )
    return findings


def validate_case_study_deck(
    deck_path: Path,
    request: CaseStudyRequest,
    evidence_paths: list[str],
) -> ValidationResult:
    policy_path = Path(
        os.getenv("TEMPLATE_POLICY_PATH", "assets/templates/contoso-template-policy.json")
    )
    result = validate_presentation(deck_path=deck_path, policy_path=policy_path, request=request)
    findings = list(result.findings)
    if result.approved:
        findings.extend(_visual_brand_findings(deck_path))
    if not evidence_paths:
        findings.append(
            ValidationFinding(
                code=FindingCode.INCONCLUSIVE,
                severity=FindingSeverity.ERROR,
                message="Evidence paths are required for source-grounded validation",
            )
        )
    elif result.approved:
        try:
            evidence = "\n".join(extract_source_text(path) for path in evidence_paths)
        except (FileNotFoundError, ValueError) as exc:
            findings.append(
                ValidationFinding(
                    code=FindingCode.INCONCLUSIVE,
                    severity=FindingSeverity.ERROR,
                    message=f"Source-grounding validation could not read the uploaded evidence: {exc}",
                )
            )
        else:
            presentation = Presentation(deck_path)
            for slide_number, slide in enumerate(presentation.slides, start=1):
                for shape in slide.shapes:
                    if not shape.name.startswith("editable:") or not getattr(shape, "has_text_frame", False):
                        continue
                    for value in _editable_content_segments(shape):
                        if _is_template_boilerplate(shape.name, value, request):
                            continue
                        if value == EVIDENCE_PENDING or text_is_evidence_backed(
                            _source_form(value, request),
                            evidence,
                        ):
                            continue
                        findings.append(
                            ValidationFinding(
                                code=FindingCode.UNSUPPORTED_EVIDENCE,
                                severity=FindingSeverity.ERROR,
                                message="Deck contains content not supported by uploaded evidence",
                                slide_number=slide_number,
                                shape_name=shape.name,
                            )
                        )
    return ValidationResult(
        approved=not findings,
        findings=tuple(findings),
        policy_version=result.policy_version,
    )


@tool(
    name="validate_case_study_deck",
    description=(
        "Deterministically validate a session PPTX against the canonical template, safety policy, "
        "source evidence paths, and approved visual brand layout."
    ),
    approval_mode="never_require",
)
def validate_case_study_deck_tool(
    relative_deck_path: str,
    request_json: str,
    source_paths_json: str,
) -> str:
    request = CaseStudyRequest.model_validate_json(request_json)
    source_paths = json.loads(source_paths_json)
    if not isinstance(source_paths, list) or not all(isinstance(path, str) for path in source_paths):
        raise ValueError("source_paths_json must be a JSON array of session workspace paths")
    result = validate_case_study_deck(
        workspace_deck_path(relative_deck_path),
        request,
        source_paths,
    )
    return result.model_dump_json()
