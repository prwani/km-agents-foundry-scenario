from __future__ import annotations

import json
import os
import pathlib
import shutil
import tempfile
import unittest
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Pt

from km_agents.agents.hosted.case_study_generator import operations
from km_agents.agents.hosted.case_study_generator.validation import (
    validate_case_study_deck,
)
from km_agents.contracts import CaseStudyContent, CaseStudyRequest, FindingCode


ROOT = pathlib.Path(__file__).resolve().parents[1]


def _presentation_text(presentation: Presentation) -> str:
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


class HostedOperationsTests(unittest.TestCase):
    def test_extracts_docx_and_generates_validated_deck_in_session_workspace(self):
        with tempfile.TemporaryDirectory() as workspace:
            workspace_path = pathlib.Path(workspace)
            input_path = workspace_path / "input"
            input_path.mkdir()
            shutil.copy(
                ROOT / "evaluation" / "corpus" / "v1" / "sources" / "fabrikam-clean-brief.docx",
                input_path / "brief.docx",
            )

            with patch.dict(os.environ, {"AGENT_WORKSPACE_ROOT": workspace}, clear=False):
                extracted = operations.extract_source_text("brief.docx")
                self.assertIn("Fabrikam", extracted)
                self.assertEqual(extracted, operations.extract_source_text("input/brief.docx"))

                content = CaseStudyContent(
                    customer_display_name="Customer",
                    title="Modernizing a customer service platform",
                    challenge="Legacy service operations limited visibility and agility.",
                    solution_overview="A phased modernization program aligned teams and data.",
                    architecture_components=["Experience", "Intelligence", "Data"],
                    implementation_steps=["Assess", "Design", "Deliver", "Measure"],
                    measurable_outcomes=["Improved visibility", "Faster delivery"],
                    customer_quote="The program created a foundation for responsible scale.",
                    next_steps=["Confirm scope", "Measure adoption"],
                    provenance_files=["brief.docx"],
                )
                result = operations.create_case_study_deck_tool(
                    content.model_dump(),
                    "case-study.pptx",
                    CaseStudyRequest(
                        customer_name="Fabrikam",
                        opportunity_summary="Modernize customer service operations.",
                        audience="Executive sponsors",
                        correlation_id="hosted-operations-test",
                    ).model_dump_json(),
                )

                self.assertIn("output/case-study.pptx", result)
                deck_path = workspace_path / "output" / "case-study.pptx"
                self.assertTrue(deck_path.is_file())
                validation = validate_case_study_deck(
                    deck_path=deck_path,
                    request=CaseStudyRequest(
                        customer_name="Fabrikam",
                        opportunity_summary="Modernize customer service operations.",
                        audience="Executive sponsors",
                        correlation_id="hosted-operations-test",
                    ),
                    evidence_paths=["brief.docx"],
                )
                self.assertTrue(validation.approved, validation.findings)

    def test_replaces_unsupported_outcomes_with_evidence_pending(self):
        with tempfile.TemporaryDirectory() as workspace:
            workspace_path = pathlib.Path(workspace)
            input_path = workspace_path / "input"
            input_path.mkdir()
            shutil.copy(
                ROOT / "evaluation" / "corpus" / "v1" / "sources" / "fabrikam-clean-brief.docx",
                input_path / "brief.docx",
            )
            content = CaseStudyContent(
                customer_display_name="Customer",
                title="Fabrikam service operations modernization",
                challenge=(
                    "Support specialists spent 42 minutes on average triaging complex cases "
                    "across disconnected knowledge sources."
                ),
                solution_overview=(
                    "A governed KM portal routes validated user uploads to Foundry agents in the "
                    "signed-in user's authorization context."
                ),
                architecture_components=["Governed KM portal", "Validated user uploads", "Foundry agents"],
                implementation_steps=[
                    "A governed KM portal routes validated user uploads to Foundry agents in the signed-in user's authorization context.",
                    "A governed KM portal routes validated user uploads to Foundry agents in the signed-in user's authorization context.",
                    "A governed KM portal routes validated user uploads to Foundry agents in the signed-in user's authorization context.",
                ],
                measurable_outcomes=[
                    "Synthetic telemetry recorded a 35 percent reduction in triage time.",
                    "Policy checks completed reached 100 percent.",
                ],
                customer_quote="Synthetic evaluation evidence - no real customer data.",
                next_steps=["Evidence pending", "Evidence pending"],
                provenance_files=["brief.docx"],
            )
            with patch.dict(os.environ, {"AGENT_WORKSPACE_ROOT": workspace}, clear=False):
                result = json.loads(
                    operations.create_case_study_deck_tool(
                        content.model_dump(),
                        "case-study.pptx",
                        CaseStudyRequest(
                            customer_name="Fabrikam",
                            opportunity_summary="Modernize customer service operations.",
                            audience="Executive sponsors",
                            correlation_id="hosted-operations-test",
                        ).model_dump_json(),
                    )
                )
                deck = Presentation(workspace_path / "output" / "case-study.pptx")

            outcome_text = "\n".join(
                shape.text
                for shape in deck.slides[5].shapes
                if getattr(shape, "has_text_frame", False)
            )
            self.assertEqual(result["suppressed_unsupported_content_fields"], 1)
            self.assertIn("35 percent reduction in triage time", outcome_text)
            self.assertIn("Evidence pending", outcome_text)
            self.assertNotIn("Fabrikam", _presentation_text(deck))
            for shape in deck.slides[5].shapes:
                if shape.name == "editable:s6:outcome-1":
                    shape.text = "Policy checks completed reached 100 percent."
            deck.save(workspace_path / "output" / "case-study.pptx")
            with patch.dict(os.environ, {"AGENT_WORKSPACE_ROOT": workspace}, clear=False):
                validation = validate_case_study_deck(
                    workspace_path / "output" / "case-study.pptx",
                    CaseStudyRequest(
                        customer_name="Fabrikam",
                        opportunity_summary="Modernize customer service operations.",
                        audience="Executive sponsors",
                        correlation_id="hosted-operations-test",
                    ),
                    ["brief.docx"],
                )
            self.assertFalse(validation.approved)
            self.assertIn(FindingCode.UNSUPPORTED_EVIDENCE, {finding.code for finding in validation.findings})

    def test_validation_rejects_editable_text_that_breaks_brand_typography(self):
        with tempfile.TemporaryDirectory() as workspace:
            workspace_path = pathlib.Path(workspace)
            input_path = workspace_path / "input"
            input_path.mkdir()
            shutil.copy(
                ROOT / "evaluation" / "corpus" / "v1" / "sources" / "fabrikam-clean-brief.docx",
                input_path / "brief.docx",
            )
            request = CaseStudyRequest(
                customer_name="Fabrikam",
                opportunity_summary="Modernize customer service operations.",
                audience="Executive sponsors",
                correlation_id="hosted-visual-qa-test",
            )
            content = CaseStudyContent(
                customer_display_name="Customer",
                title="Modernizing a customer service platform",
                challenge="Legacy service operations limited visibility and agility.",
                solution_overview="A phased modernization program aligned teams and data.",
                architecture_components=["Experience", "Intelligence", "Data"],
                implementation_steps=["Assess", "Design", "Deliver", "Measure"],
                measurable_outcomes=["Improved visibility", "Faster delivery"],
                customer_quote="The program created a foundation for responsible scale.",
                next_steps=["Confirm scope", "Measure adoption"],
                provenance_files=["brief.docx"],
            )
            with patch.dict(os.environ, {"AGENT_WORKSPACE_ROOT": workspace}, clear=False):
                operations.create_case_study_deck_tool(
                    content.model_dump(),
                    "case-study.pptx",
                    request.model_dump_json(),
                )
                deck_path = workspace_path / "output" / "case-study.pptx"
                deck = Presentation(deck_path)
                title = next(
                    shape for shape in deck.slides[0].shapes if shape.name == "editable:s1:title"
                )
                title.text_frame.paragraphs[0].runs[0].font.size = Pt(72)
                deck.save(deck_path)
                validation = validate_case_study_deck(deck_path, request, ["brief.docx"])

            self.assertFalse(validation.approved)
            self.assertIn(
                FindingCode.VISUAL_BRAND_VIOLATION,
                {finding.code for finding in validation.findings},
            )

    def test_generation_tool_accepts_incomplete_evidence_lists_and_pads_template_slots(self):
        content = {
            "customer_display_name": "Customer",
            "title": "Fabrikam service operations modernization",
            "challenge": "Support specialists spent 42 minutes on average triaging complex cases across disconnected knowledge sources.",
            "solution_overview": "A governed KM portal routes validated user uploads to Foundry agents in the signed-in user's authorization context.",
            "architecture_components": ["A governed KM portal"],
            "implementation_steps": ["Evidence pending"],
            "measurable_outcomes": ["Synthetic telemetry recorded a 35 percent reduction in triage time and 22 percent fewer handoffs."],
            "customer_quote": "Evidence pending",
            "next_steps": ["Evidence pending"],
            "provenance_files": ["input/fabrikam-clean-brief.docx"],
        }
        request_json = CaseStudyRequest(
            customer_name="Fabrikam",
            opportunity_summary="Modernize customer service operations.",
            audience="Executive sponsors",
            correlation_id="hosted-operations-test",
        ).model_dump_json()

        parsed = operations.create_case_study_deck_tool.input_model.model_validate(
            {
                "content": content,
                "output_filename": "case-study.pptx",
                "request_json": request_json,
            }
        )

        self.assertEqual(parsed.content.measurable_outcomes, content["measurable_outcomes"])
        self.assertEqual(parsed.content.next_steps, content["next_steps"])


if __name__ == "__main__":
    unittest.main()
