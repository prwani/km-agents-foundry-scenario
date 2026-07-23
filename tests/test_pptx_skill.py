import json
import pathlib
import tempfile
import unittest

from pptx import Presentation

from km_agents.contracts import (
    CaseStudyContent,
    CaseStudyRequest,
    ImplementationKind,
)
from km_agents.pptx_skill.generation import generate_case_study
from km_agents.pptx_skill.validation import validate_presentation


ROOT = pathlib.Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "assets" / "templates" / "contoso-case-study-template.pptx"
POLICY = ROOT / "assets" / "templates" / "contoso-template-policy.json"


def make_request(approved_name: bool = True) -> CaseStudyRequest:
    return CaseStudyRequest(
        implementation=ImplementationKind.HOSTED,
        customer_name="Fabrikam",
        customer_name_approved_for_external_use=approved_name,
        opportunity_summary="Improve service operations",
        audience="executives",
        correlation_id="pptx-test-001",
    )


def make_content(customer_name: str = "Fabrikam") -> CaseStudyContent:
    return CaseStudyContent(
        customer_display_name=customer_name,
        title="Transforming service operations with responsible AI",
        challenge="Teams needed faster access to trusted knowledge without weakening governance.",
        solution_overview="Contoso introduced a governed agent experience grounded in approved enterprise content.",
        architecture_components=[
            "KM portal",
            "Foundry agents",
            "Direct file upload",
            "Foundry OBO",
            "Validation policy",
            "Application Insights",
        ],
        implementation_steps=["Discover", "Design", "Pilot", "Scale"],
        measurable_outcomes=["35% faster triage", "22% fewer handoffs", "100% policy checks"],
        customer_quote="The governed workflow helped our teams move faster with confidence.",
        next_steps=["Expand the synthetic corpus", "Measure adoption", "Review policy quarterly"],
        provenance_files=["brief.docx"],
    )


class PptxSkillTests(unittest.TestCase):
    def test_template_and_policy_are_versioned(self):
        self.assertTrue(TEMPLATE.is_file())
        policy = json.loads(POLICY.read_text(encoding="utf-8"))
        self.assertEqual(policy["slide_count"], 8)
        self.assertEqual(policy["policy_version"], "1.0.0")
        self.assertEqual(len(policy["template_sha256"]), 64)

    def test_generated_deck_passes_deterministic_validation(self):
        with tempfile.TemporaryDirectory() as directory:
            output = pathlib.Path(directory) / "case-study.pptx"
            generate_case_study(TEMPLATE, output, make_content())
            result = validate_presentation(output, POLICY, make_request())
        self.assertTrue(result.approved, result.reasons)

    def test_generated_deck_preserves_template_text_colors(self):
        with tempfile.TemporaryDirectory() as directory:
            output = pathlib.Path(directory) / "case-study.pptx"
            generate_case_study(TEMPLATE, output, make_content())
            template = Presentation(TEMPLATE)
            generated = Presentation(output)

        for slide_index, shape_name in (
            (0, "editable:s1:title"),
            (1, "editable:s2:context"),
            (6, "editable:s7:quote"),
            (7, "editable:s8:next-steps"),
        ):
            template_shape = next(
                shape for shape in template.slides[slide_index].shapes if shape.name == shape_name
            )
            generated_shape = next(
                shape for shape in generated.slides[slide_index].shapes if shape.name == shape_name
            )
            self.assertEqual(
                generated_shape.text_frame.paragraphs[0].runs[0].font.color.rgb,
                template_shape.text_frame.paragraphs[0].runs[0].font.color.rgb,
                shape_name,
            )

    def test_protected_shape_change_is_rejected(self):
        with tempfile.TemporaryDirectory() as directory:
            output = pathlib.Path(directory) / "case-study.pptx"
            generate_case_study(TEMPLATE, output, make_content())
            deck = Presentation(output)
            protected = next(
                shape
                for shape in deck.slides[0].shapes
                if shape.name == "protected:s1:logo-wordmark"
            )
            protected.text = "ALTERED BRAND"
            deck.save(output)
            result = validate_presentation(output, POLICY, make_request())
        self.assertFalse(result.approved)
        self.assertIn("protected element", " ".join(result.reasons).lower())

    def test_unapproved_customer_name_is_rejected(self):
        with tempfile.TemporaryDirectory() as directory:
            output = pathlib.Path(directory) / "case-study.pptx"
            generate_case_study(TEMPLATE, output, make_content("Fabrikam"))
            result = validate_presentation(output, POLICY, make_request(approved_name=False))
        self.assertFalse(result.approved)
        self.assertIn("without external-use attestation", " ".join(result.reasons))


if __name__ == "__main__":
    unittest.main()
