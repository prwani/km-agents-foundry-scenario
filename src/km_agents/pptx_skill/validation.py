from __future__ import annotations

import json
from pathlib import Path
import zipfile
from typing import Any

from pptx import Presentation

from km_agents.contracts import (
    CaseStudyRequest,
    FindingCode,
    FindingSeverity,
    ValidationFinding,
    ValidationResult,
)
from km_agents.safety import find_sensitive_markers


def _shape_text(shape: Any) -> str:
    return shape.text if getattr(shape, "has_text_frame", False) else ""


def _fingerprint(shape: Any) -> dict[str, Any]:
    return {
        "shape_type": int(shape.shape_type),
        "left": int(shape.left),
        "top": int(shape.top),
        "width": int(shape.width),
        "height": int(shape.height),
        "text": _shape_text(shape),
    }


def _error(
    code: FindingCode,
    message: str,
    *,
    slide_number: int | None = None,
    shape_name: str | None = None,
    evidence: dict[str, Any] | None = None,
) -> ValidationFinding:
    return ValidationFinding(
        code=code,
        severity=FindingSeverity.ERROR,
        message=message,
        slide_number=slide_number,
        shape_name=shape_name,
        evidence=evidence or {},
    )


def validate_presentation(
    deck_path: Path,
    policy_path: Path,
    request: CaseStudyRequest,
) -> ValidationResult:
    if not deck_path.is_file() or not zipfile.is_zipfile(deck_path):
        return ValidationResult(
            approved=False,
            policy_version="unknown",
            findings=(_error(FindingCode.INVALID_FILE, "Deck is missing or not a valid PPTX package"),),
        )
    if not policy_path.is_file():
        return ValidationResult(
            approved=False,
            policy_version="unknown",
            findings=(_error(FindingCode.INCONCLUSIVE, "Template policy is unavailable"),),
        )

    policy = json.loads(policy_path.read_text(encoding="utf-8"))
    policy_version = str(policy["policy_version"])
    presentation = Presentation(deck_path)
    findings: list[ValidationFinding] = []

    expected_count = int(policy["slide_count"])
    if len(presentation.slides) != expected_count:
        findings.append(
            _error(
                FindingCode.TEMPLATE_MISMATCH,
                f"Expected {expected_count} slides, found {len(presentation.slides)}",
            )
        )

    for slide_policy in policy["slides"]:
        slide_number = int(slide_policy["slide_number"])
        if slide_number > len(presentation.slides):
            continue
        slide = presentation.slides[slide_number - 1]
        actual = {shape.name: shape for shape in slide.shapes}
        for shape_name in slide_policy["editable_shapes"]:
            if shape_name not in actual:
                findings.append(
                    _error(
                        FindingCode.REQUIRED_CONTENT_MISSING,
                        "Required editable region is missing",
                        slide_number=slide_number,
                        shape_name=shape_name,
                    )
                )
        for shape_name, expected in slide_policy["protected_shapes"].items():
            shape = actual.get(shape_name)
            if shape is None:
                findings.append(
                    _error(
                        FindingCode.PROTECTED_ELEMENT_CHANGED,
                        "Protected element is missing",
                        slide_number=slide_number,
                        shape_name=shape_name,
                    )
                )
                continue
            actual_fingerprint = _fingerprint(shape)
            if actual_fingerprint != expected:
                findings.append(
                    _error(
                        FindingCode.PROTECTED_ELEMENT_CHANGED,
                        "Protected element differs from the canonical template",
                        slide_number=slide_number,
                        shape_name=shape_name,
                        evidence={"expected": expected, "actual": actual_fingerprint},
                    )
                )

    all_text = "\n".join(
        _shape_text(shape)
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    markers = find_sensitive_markers(all_text)
    if markers:
        findings.append(
            _error(
                FindingCode.SENSITIVE_INFORMATION,
                "Deck contains potential business-sensitive information",
                evidence={"marker_count": len(markers)},
            )
        )
    if (
        not request.customer_name_approved_for_external_use
        and request.customer_name.casefold() in all_text.casefold()
    ):
        findings.append(
            _error(
                FindingCode.CUSTOMER_NAME_NOT_APPROVED,
                "Customer name appears without external-use attestation",
            )
        )

    return ValidationResult(
        approved=not findings,
        findings=tuple(findings),
        policy_version=policy_version,
    )
