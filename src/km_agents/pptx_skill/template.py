from __future__ import annotations

import argparse
import hashlib
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .branding import CONTOSO


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _set_name(shape: Any, name: str) -> None:
    shape._element.nvSpPr.cNvPr.set("name", name)


def _text(
    slide: Any,
    name: str,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    color: str = CONTOSO.ink,
    bold: bool = False,
    font: str = CONTOSO.body_font,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set_name(shape, name)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return shape


def _rect(
    slide: Any,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    color: str,
    *,
    radius: bool = False,
    line_color: str | None = None,
) -> Any:
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    _set_name(shape, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.color.rgb = _rgb(line_color or color)
    return shape


def _circle(
    slide: Any,
    name: str,
    x: float,
    y: float,
    diameter: float,
    color: str,
) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(diameter),
        Inches(diameter),
    )
    _set_name(shape, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.color.rgb = _rgb(color)
    return shape


def _add_logo(slide: Any, slide_number: int, *, inverse: bool = False) -> None:
    mark_color = CONTOSO.gold if inverse else CONTOSO.teal
    text_color = CONTOSO.white if inverse else CONTOSO.navy
    _circle(
        slide,
        f"protected:s{slide_number}:logo-mark",
        0.55,
        0.38,
        0.45,
        mark_color,
    )
    _text(
        slide,
        f"protected:s{slide_number}:logo-initials",
        "CL",
        0.64,
        0.48,
        0.28,
        0.18,
        size=10,
        color=CONTOSO.navy,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    _text(
        slide,
        f"protected:s{slide_number}:logo-wordmark",
        "CONTOSO LIMITED",
        1.12,
        0.43,
        2.25,
        0.28,
        size=12,
        color=text_color,
        bold=True,
        font=CONTOSO.header_font,
    )


def _add_footer(slide: Any, slide_number: int, *, inverse: bool = False) -> None:
    color = CONTOSO.white if inverse else CONTOSO.muted
    _text(
        slide,
        f"protected:s{slide_number}:confidentiality",
        "SYNTHETIC DEMO • CONTOSO LIMITED",
        0.55,
        7.08,
        3.25,
        0.2,
        size=8,
        color=color,
        bold=True,
    )
    _text(
        slide,
        f"protected:s{slide_number}:slide-number",
        f"{slide_number:02d}",
        12.2,
        7.04,
        0.55,
        0.22,
        size=9,
        color=color,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def _add_standard_chrome(slide: Any, slide_number: int, title: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(CONTOSO.cloud)
    _add_logo(slide, slide_number)
    _text(
        slide,
        f"protected:s{slide_number}:title",
        title,
        0.65,
        1.08,
        11.95,
        0.58,
        size=30,
        color=CONTOSO.navy,
        bold=True,
        font=CONTOSO.header_font,
    )
    _add_footer(slide, slide_number)


def _slide_title(prs: Presentation) -> None:
    slide_number = 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(CONTOSO.navy)
    _add_logo(slide, slide_number, inverse=True)
    _rect(
        slide,
        "protected:s1:accent-panel",
        9.65,
        0,
        3.683,
        7.5,
        CONTOSO.teal,
    )
    _circle(slide, "protected:s1:accent-orbit-1", 10.25, 1.35, 1.55, CONTOSO.gold)
    _circle(slide, "protected:s1:accent-orbit-2", 11.18, 4.2, 0.78, CONTOSO.white)
    _text(
        slide,
        "editable:s1:customer-name",
        "{{CUSTOMER_NAME}}",
        0.75,
        1.55,
        7.9,
        0.45,
        size=18,
        color=CONTOSO.gold,
        bold=True,
    )
    _text(
        slide,
        "editable:s1:title",
        "{{TITLE}}",
        0.75,
        2.12,
        8.15,
        1.45,
        size=40,
        color=CONTOSO.white,
        bold=True,
        font=CONTOSO.header_font,
        valign=MSO_ANCHOR.MIDDLE,
    )
    _text(
        slide,
        "editable:s1:subtitle",
        "A synthetic customer success story",
        0.75,
        3.83,
        7.5,
        0.45,
        size=18,
        color="CAD7E3",
    )
    _add_footer(slide, slide_number, inverse=True)


def _slide_challenge(prs: Presentation) -> None:
    slide_number = 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_standard_chrome(slide, slide_number, "The customer challenge")
    _rect(
        slide,
        "protected:s2:challenge-card",
        0.7,
        1.95,
        7.15,
        4.5,
        CONTOSO.white,
        radius=True,
        line_color="D7E1EA",
    )
    _rect(slide, "protected:s2:context-panel", 8.25, 1.95, 4.35, 4.5, CONTOSO.navy, radius=True)
    _text(
        slide,
        "protected:s2:challenge-label",
        "THE CHALLENGE",
        1.05,
        2.45,
        3.45,
        0.35,
        size=13,
        color=CONTOSO.teal,
        bold=True,
    )
    _text(
        slide,
        "editable:s2:challenge",
        "{{CHALLENGE}}",
        1.05,
        3.0,
        6.45,
        2.45,
        size=22,
        color=CONTOSO.ink,
        valign=MSO_ANCHOR.MIDDLE,
    )
    _text(
        slide,
        "protected:s2:context-label",
        "WHY NOW",
        8.7,
        2.45,
        3.45,
        0.35,
        size=13,
        color=CONTOSO.gold,
        bold=True,
    )
    _text(
        slide,
        "editable:s2:context",
        "{{OPPORTUNITY_CONTEXT}}",
        8.7,
        3.0,
        3.25,
        2.5,
        size=18,
        color=CONTOSO.white,
        valign=MSO_ANCHOR.MIDDLE,
    )


def _slide_solution(prs: Presentation) -> None:
    slide_number = 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_standard_chrome(slide, slide_number, "Solution overview")
    _text(
        slide,
        "editable:s3:solution",
        "{{SOLUTION_OVERVIEW}}",
        0.75,
        2.28,
        5.2,
        2.65,
        size=22,
        color=CONTOSO.ink,
        valign=MSO_ANCHOR.MIDDLE,
    )
    for index, (label, color) in enumerate(
        (("DISCOVER", CONTOSO.navy), ("DESIGN", CONTOSO.teal), ("DELIVER", CONTOSO.gold))
    ):
        x = 6.4 + index * 2.05
        _circle(slide, f"protected:s3:step-{index + 1}", x, 2.25, 1.25, color)
        _text(
            slide,
            f"protected:s3:step-icon-{index + 1}",
            str(index + 1),
            x + 0.34,
            2.64,
            0.56,
            0.3,
            size=18,
            color=CONTOSO.white if index < 2 else CONTOSO.navy,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _text(
            slide,
            f"protected:s3:step-label-{index + 1}",
            label,
            x - 0.15,
            3.72,
            1.55,
            0.3,
            size=12,
            color=CONTOSO.navy,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    _text(
        slide,
        "editable:s3:solution-pillars",
        "{{SOLUTION_PILLARS}}",
        6.35,
        4.45,
        5.75,
        1.15,
        size=15,
        color=CONTOSO.muted,
        align=PP_ALIGN.CENTER,
    )


def _slide_architecture(prs: Presentation) -> None:
    slide_number = 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_standard_chrome(slide, slide_number, "Architecture")
    _rect(slide, "protected:s4:architecture-canvas", 0.72, 1.9, 11.88, 4.75, CONTOSO.white, radius=True)
    columns = (
        ("EXPERIENCE", 1.1, CONTOSO.navy),
        ("INTELLIGENCE", 4.75, CONTOSO.teal),
        ("DATA & ACTIONS", 8.4, CONTOSO.gold),
    )
    for index, (label, x, color) in enumerate(columns, start=1):
        _rect(slide, f"protected:s4:column-{index}", x, 2.45, 3.0, 3.45, "F8FAFC", radius=True, line_color=color)
        _text(
            slide,
            f"protected:s4:column-label-{index}",
            label,
            x + 0.25,
            2.75,
            2.5,
            0.3,
            size=12,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _text(
            slide,
            f"editable:s4:architecture-column-{index}",
            f"{{{{ARCHITECTURE_COLUMN_{index}}}}}",
            x + 0.28,
            3.3,
            2.44,
            1.9,
            size=15,
            color=CONTOSO.ink,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )


def _slide_journey(prs: Presentation) -> None:
    slide_number = 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_standard_chrome(slide, slide_number, "Implementation journey")
    _rect(slide, "protected:s5:timeline-track", 1.675, 3.42, 9.0, 0.12, CONTOSO.teal)
    for index in range(4):
        x = 1.25 + index * 3.0
        _circle(slide, f"protected:s5:milestone-{index + 1}", x, 3.05, 0.85, CONTOSO.navy if index % 2 == 0 else CONTOSO.teal)
        _text(
            slide,
            f"protected:s5:milestone-number-{index + 1}",
            str(index + 1),
            x + 0.2,
            3.27,
            0.45,
            0.22,
            size=12,
            color=CONTOSO.white,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _text(
            slide,
            f"editable:s5:implementation-step-{index + 1}",
            f"{{{{IMPLEMENTATION_STEP_{index + 1}}}}}",
            x - 0.55,
            4.25,
            1.95,
            1.05,
            size=15,
            color=CONTOSO.ink,
            bold=True,
            align=PP_ALIGN.CENTER,
        )


def _slide_outcomes(prs: Presentation) -> None:
    slide_number = 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_standard_chrome(slide, slide_number, "Measurable outcomes")
    for index, color in enumerate((CONTOSO.navy, CONTOSO.teal, CONTOSO.gold)):
        x = 0.85 + index * 4.1
        _rect(slide, f"protected:s6:metric-card-{index + 1}", x, 2.1, 3.65, 3.75, CONTOSO.white, radius=True)
        _circle(slide, f"protected:s6:metric-icon-{index + 1}", x + 1.35, 2.55, 0.95, color)
        _text(
            slide,
            f"protected:s6:metric-number-{index + 1}",
            str(index + 1),
            x + 1.61,
            2.86,
            0.42,
            0.24,
            size=13,
            color=CONTOSO.white if index < 2 else CONTOSO.navy,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _text(
            slide,
            f"editable:s6:outcome-{index + 1}",
            f"{{{{MEASURABLE_OUTCOME_{index + 1}}}}}",
            x + 0.35,
            3.85,
            2.95,
            1.3,
            size=17,
            color=CONTOSO.ink,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    _text(
        slide,
        "protected:s6:evidence-note",
        "Only source-supported outcomes are permitted",
        3.45,
        6.25,
        6.45,
        0.28,
        size=11,
        color=CONTOSO.muted,
        align=PP_ALIGN.CENTER,
    )


def _slide_quote(prs: Presentation) -> None:
    slide_number = 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(CONTOSO.teal)
    _add_logo(slide, slide_number, inverse=True)
    _text(
        slide,
        "protected:s7:quote-mark",
        "“",
        0.72,
        1.38,
        1.1,
        1.15,
        size=72,
        color=CONTOSO.gold,
        bold=True,
        font="Georgia",
    )
    _text(
        slide,
        "editable:s7:quote",
        "{{CUSTOMER_QUOTE}}",
        1.6,
        2.0,
        10.65,
        2.45,
        size=30,
        color=CONTOSO.white,
        font="Georgia",
        valign=MSO_ANCHOR.MIDDLE,
    )
    _text(
        slide,
        "editable:s7:quote-attribution",
        "{{QUOTE_ATTRIBUTION}}",
        1.65,
        5.0,
        10.55,
        0.4,
        size=15,
        color="D7FFFF",
        bold=True,
    )
    _add_footer(slide, slide_number, inverse=True)


def _slide_next_steps(prs: Presentation) -> None:
    slide_number = 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(CONTOSO.navy)
    _add_logo(slide, slide_number, inverse=True)
    _text(
        slide,
        "protected:s8:title",
        "Next steps",
        0.75,
        1.22,
        5.2,
        0.65,
        size=36,
        color=CONTOSO.white,
        bold=True,
        font=CONTOSO.header_font,
    )
    _text(
        slide,
        "editable:s8:next-steps",
        "{{NEXT_STEPS}}",
        0.8,
        2.18,
        6.85,
        3.1,
        size=22,
        color=CONTOSO.white,
        valign=MSO_ANCHOR.MIDDLE,
    )
    _rect(slide, "protected:s8:closing-panel", 8.35, 1.45, 4.1, 4.95, CONTOSO.teal, radius=True)
    _text(
        slide,
        "protected:s8:closing-message",
        "BUILD TRUST.\nPROVE VALUE.\nSCALE RESPONSIBLY.",
        8.75,
        2.65,
        3.3,
        1.85,
        size=23,
        color=CONTOSO.white,
        bold=True,
        font=CONTOSO.header_font,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    _add_footer(slide, slide_number, inverse=True)


def _shape_text(shape: Any) -> str:
    return shape.text if getattr(shape, "has_text_frame", False) else ""


def _shape_fingerprint(shape: Any) -> dict[str, Any]:
    return {
        "shape_type": int(shape.shape_type),
        "left": int(shape.left),
        "top": int(shape.top),
        "width": int(shape.width),
        "height": int(shape.height),
        "text": _shape_text(shape),
    }


def build_policy(template_path: Path) -> dict[str, Any]:
    presentation = Presentation(template_path)
    slides: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        editable: list[str] = []
        protected: dict[str, dict[str, Any]] = {}
        for shape in slide.shapes:
            if shape.name.startswith("editable:"):
                editable.append(shape.name)
            elif shape.name.startswith("protected:"):
                protected[shape.name] = _shape_fingerprint(shape)
        slides.append(
            {
                "slide_number": slide_number,
                "editable_shapes": sorted(editable),
                "protected_shapes": protected,
            }
        )
    return {
        "policy_version": "1.0.0",
        "template_sha256": hashlib.sha256(template_path.read_bytes()).hexdigest(),
        "slide_count": len(presentation.slides),
        "slides": slides,
        "customer_name_policy": "anonymize_unless_attested",
        "allowed_extensions": ["docx", "pptx", "pdf", "xlsx"],
    }


def create_template(template_path: Path, policy_path: Path) -> None:
    template_path.parent.mkdir(parents=True, exist_ok=True)
    policy_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    presentation.core_properties.title = "Contoso Limited Case Study Template"
    presentation.core_properties.subject = "Synthetic case-study template"
    presentation.core_properties.author = "Contoso Limited"
    presentation.core_properties.keywords = "synthetic,case-study,template"
    _slide_title(presentation)
    _slide_challenge(presentation)
    _slide_solution(presentation)
    _slide_architecture(presentation)
    _slide_journey(presentation)
    _slide_outcomes(presentation)
    _slide_quote(presentation)
    _slide_next_steps(presentation)
    presentation.save(template_path)
    policy_path.write_text(json.dumps(build_policy(template_path), indent=2), encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate the Contoso Limited case-study template.")
    parser.add_argument(
        "--template",
        type=Path,
        default=Path("assets/templates/contoso-case-study-template.pptx"),
    )
    parser.add_argument(
        "--policy",
        type=Path,
        default=Path("assets/templates/contoso-template-policy.json"),
    )
    args = parser.parse_args()
    create_template(args.template, args.policy)


if __name__ == "__main__":
    main()
