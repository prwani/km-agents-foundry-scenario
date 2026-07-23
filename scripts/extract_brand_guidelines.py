"""Extract textual brand guidance from the Contoso brand-guidelines reference deck.

The Hosted case-study generator cannot open the brand-guidelines PPTX directly: its
file-access tool is sandboxed to the agent workspace and the guidelines asset lives
outside it. Instead, the generator reads a small extracted markdown summary that is
embedded directly into the agent's instructions at startup.

This script is the *only* supported way to produce that summary. Whenever the brand
deck at ``assets/templates/contoso-case-study-template-with-brand-guidelines.pptx``
changes, re-run this script and commit the regenerated markdown alongside the updated
deck:

    python scripts/extract_brand_guidelines.py

It writes the same content to both the root and ``src`` copies of
``assets/templates/contoso-brand-guidelines.md`` (the project keeps duplicate asset
copies in sync, matching the existing PPTX template convention) unless ``--output``
is given explicitly.

The output file embeds the source deck's SHA-256 hash in a header comment. Run
``python scripts/extract_brand_guidelines.py --check`` in CI or before a release to
confirm the committed markdown still matches the current brand deck; it exits
non-zero if the markdown is stale and needs regenerating.
"""

from __future__ import annotations

import argparse
import hashlib
import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_SOURCE = (
    ROOT / "assets" / "templates" / "contoso-case-study-template-with-brand-guidelines.pptx"
)
DEFAULT_OUTPUTS = (
    ROOT / "assets" / "templates" / "contoso-brand-guidelines.md",
    ROOT / "src" / "assets" / "templates" / "contoso-brand-guidelines.md",
)
# Slides 1-8 are the canonical case-study template; guidance content starts after them.
CANONICAL_SLIDE_COUNT = 8


def _slide_text_lines(slide) -> list[str]:
    lines: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if text:
                lines.append(text)
    return lines


def _drop_boilerplate(all_slide_lines: list[list[str]]) -> list[list[str]]:
    """Remove running headers/footers/page numbers repeated across most slides."""
    counts: Counter[str] = Counter()
    for lines in all_slide_lines:
        for line in set(lines):
            counts[line] += 1
    total = len(all_slide_lines)
    boilerplate = {line for line, count in counts.items() if total and count / total >= 0.5}
    cleaned = []
    for lines in all_slide_lines:
        kept = [
            line
            for line in lines
            if line not in boilerplate and not (line.isdigit() and len(line) <= 2)
        ]
        cleaned.append(kept)
    return cleaned


def extract_markdown(source: Path, skip_slides: int = CANONICAL_SLIDE_COUNT) -> str:
    presentation = Presentation(str(source))
    slides = list(presentation.slides)
    guidance_slides = slides[skip_slides:]
    if not guidance_slides:
        raise ValueError(
            f"No guidance slides found after skipping the first {skip_slides} canonical "
            f"slides in {source}"
        )

    raw_lines = [_slide_text_lines(slide) for slide in guidance_slides]
    cleaned_lines = _drop_boilerplate(raw_lines)
    source_hash = hashlib.sha256(source.read_bytes()).hexdigest()

    parts = [
        "<!-- AUTO-GENERATED FILE. Do not edit by hand. -->",
        "<!--",
        "  Regenerate with: python scripts/extract_brand_guidelines.py",
        f"  Source deck: {source.relative_to(ROOT).as_posix()}",
        f"  source_sha256: {source_hash}",
        "-->",
        "",
        "# Contoso Brand Guidelines (extracted reference)",
        "",
        "This file is generated from the brand-guidelines reference deck and is read by the "
        "case-study generator agent at runtime as design and voice guidance. It must never be "
        "copied into a generated customer-facing deck.",
        "",
    ]
    for index, lines in enumerate(cleaned_lines, start=skip_slides + 1):
        if not lines:
            continue
        heading = lines[0]
        parts.append(f"## Slide {index}: {heading}")
        parts.append("")
        for line in lines[1:]:
            parts.append(f"- {line}")
        parts.append("")
    return "\n".join(parts).rstrip() + "\n"


def _extract_committed_hash(markdown_path: Path) -> str | None:
    if not markdown_path.is_file():
        return None
    for line in markdown_path.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if line.startswith("source_sha256:"):
            return line.split(":", 1)[1].strip()
    return None


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("--source", type=Path, default=DEFAULT_SOURCE)
    parser.add_argument("--skip-slides", type=int, default=CANONICAL_SLIDE_COUNT)
    parser.add_argument(
        "--output",
        type=Path,
        action="append",
        default=None,
        help="Output markdown path; may be repeated. Defaults to both root and src asset copies.",
    )
    parser.add_argument(
        "--check",
        action="store_true",
        help="Do not write files; exit non-zero if any output is stale relative to --source.",
    )
    args = parser.parse_args()

    if not args.source.is_file():
        print(f"Source brand-guidelines deck not found: {args.source}", file=sys.stderr)
        return 2

    outputs = args.output or list(DEFAULT_OUTPUTS)
    markdown = extract_markdown(args.source, skip_slides=args.skip_slides)
    current_hash = hashlib.sha256(args.source.read_bytes()).hexdigest()

    if args.check:
        stale = []
        for output in outputs:
            committed_hash = _extract_committed_hash(output)
            if committed_hash != current_hash:
                stale.append(output)
        if stale:
            for output in stale:
                print(f"Stale brand-guidelines markdown: {output}", file=sys.stderr)
            print(
                "Run 'python scripts/extract_brand_guidelines.py' and commit the result.",
                file=sys.stderr,
            )
            return 1
        print("Brand-guidelines markdown is up to date.")
        return 0

    for output in outputs:
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(markdown, encoding="utf-8")
        print(f"Wrote {output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
